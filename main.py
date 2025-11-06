from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID
from datetime import date


def parse_hierarchical_text(text):
    """
    インデントされたテキストを階層構造に解析する
    :param text: 入力テキスト
    :return: 階層構造を表すリスト
    """
    lines = text.strip().split("\n")
    result = []
    stack = [(-1, result)]

    for line in lines:
        indent = len(line) - len(line.lstrip())
        content = line.strip("- ").strip()

        # 現在の行のインデントレベルに基づいてスタックを調整
        while stack and indent <= stack[-1][0]:
            stack.pop()

        if not stack:
            stack = [(-1, result)]

        parent = stack[-1][1]
        current = []
        parent.append((content, current))
        stack.append((indent, current))

    return result


def add_body(text_frame, items, level=0):
    """
    テキストフレームに階層構造のテキストを追加する
    :param text_frame: PowerPointのテキストフレーム
    :param items: 階層構造のアイテムリスト
    :param level: 現在のインデントレベル
    """
    for item, sub_items in items:
        # 既存のパラグラフを使用するか、新しいパラグラフを追加
        p = text_frame.add_paragraph() if text_frame.text else text_frame.paragraphs[0]
        p.text = item
        p.level = level
        # サブアイテムを再帰的に処理
        add_body(text_frame, sub_items, level + 1)


def set_japanese_language(shape):
    """
    シェイプ内のすべてのテキストの言語を日本語に設定する
    言語を適切に設定しないと、不必要な校正の赤線が引かれる
    :param shape: PowerPointのシェイプオブジェクト
    """
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.language_id = MSO_LANGUAGE_ID.JAPANESE


def create_presentation(hierarchy, template_path, output_path):
    """
    階層構造のデータからPowerPointプレゼンテーションを作成する
    :param hierarchy: 階層構造のデータ
    :param template_path: テンプレートファイルのパス
    :param output_path: 出力ファイルのパス
    """
    prs = Presentation(template_path)

    for item, sub_items in hierarchy:
        # 新しいスライドを追加
        # テンプレートによってレイアウトの種類や順序が異なる場合がある
        slide = prs.slides.add_slide(prs.slide_layouts[2])  # slide_layouts[#]を調整する
        slide.shapes.title.text = item

        # 本文を追加
        body_shape = slide.shapes.placeholders[1]
        add_body(body_shape.text_frame, sub_items)
        set_japanese_language(body_shape)

    prs.save(output_path)
    print(f"{output_path}を作成")


def main():
    input_file = "input.txt"
    template_file = "template.pptx"
    today = date.today()
    output_file = f"out/ゼミ{today}.pptx"

    # 入力ファイルを読み込む
    with open(input_file, encoding="utf-8") as f:
        text = f.read()

    # テキストを解析し、階層構造を作成
    hierarchy = parse_hierarchical_text(text)

    # プレゼンテーションを作成
    create_presentation(hierarchy, template_file, output_file)

    print("スライド番号を忘れずに")


if __name__ == "__main__":
    main()
